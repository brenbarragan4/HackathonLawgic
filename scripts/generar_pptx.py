# -*- coding: utf-8 -*-
"""Genera la presentacion PowerPoint de demo de MILEXLEGAL."""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

# ----- Paleta (verde esmeralda profundo + dorado + crema, estilo premium) -----
NEGRO   = RGBColor(0x14, 0x3A, 0x2E)  # verde esmeralda profundo (fondos oscuros + titulos)
VERDE2  = RGBColor(0x1E, 0x4D, 0x3D)  # verde para degradados
TINTA   = RGBColor(0x2B, 0x33, 0x2D)  # texto de cuerpo (verde-carbon)
GRIS    = RGBColor(0x6B, 0x71, 0x68)  # gris salvia apagado
LINEA   = RGBColor(0xE3, 0xDD, 0xD0)  # linea calida
CREMA   = RGBColor(0xFB, 0xF8, 0xF2)  # crema de fondo
SOFT    = RGBColor(0xF0, 0xEC, 0xE0)  # panel crema calido
DORADO  = RGBColor(0xB1, 0x8A, 0x4B)  # dorado rico
DORADO2 = RGBColor(0xCB, 0xA8, 0x6A)  # dorado claro para acentos
CLARO   = RGBColor(0xD9, 0xE0, 0xD6)  # texto claro sobre verde
TENUE   = RGBColor(0xA6, 0xB2, 0xA6)  # texto tenue sobre verde
BLANCO  = RGBColor(0xFF, 0xFF, 0xFF)

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]

SERIF = "Georgia"
SANS  = "Calibri"

def fondo(slide, color=CREMA):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = color

def caja(slide, x, y, w, h, color, line=None, line_w=1.0):
    sh = slide.shapes.add_shape(1, x, y, w, h)  # rectangulo
    sh.fill.solid(); sh.fill.fore_color.rgb = color
    if line is None:
        sh.line.fill.background()
    else:
        sh.line.color.rgb = line; sh.line.width = Pt(line_w)
    sh.shadow.inherit = False
    return sh

def texto(slide, x, y, w, h, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, sp_after=6):
    """runs: lista de parrafos; cada parrafo es lista de (texto, size, color, bold, font, italic)."""
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame; tf.word_wrap = True
    tf.vertical_anchor = anchor
    for i, parr in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.space_after = Pt(sp_after)
        p.line_spacing = 1.15
        for (t, size, color, bold, font, italic) in parr:
            r = p.add_run(); r.text = t
            r.font.size = Pt(size); r.font.color.rgb = color
            r.font.bold = bold; r.font.name = font; r.font.italic = italic
    return tb

def R(t, size, color=TINTA, bold=False, font=SANS, italic=False):
    return (t, size, color, bold, font, italic)

MX, TOP = Inches(0.9), Inches(0.7)
CW = SW - MX*2

def encabezado(slide, kicker, titulo):
    fondo(slide)
    caja(slide, MX, TOP, Inches(0.12), Inches(0.9), DORADO)
    texto(slide, Inches(1.15), TOP, CW, Inches(0.4),
          [[R(kicker, 12, DORADO, True, SANS)]], sp_after=2)
    texto(slide, Inches(1.15), Inches(1.05), CW, Inches(0.9),
          [[R(titulo, 30, NEGRO, True, SERIF)]])
    caja(slide, Inches(1.15), Inches(1.95), CW - Inches(0.25), Pt(1.4), LINEA)

# ============ SLIDE 1: PORTADA ============
s = prs.slides.add_slide(BLANK)
fondo(s, NEGRO)
caja(s, Inches(0.9), Inches(2.0), Inches(1.1), Pt(4), DORADO)
texto(s, Inches(0.9), Inches(2.15), CW, Inches(0.5),
      [[R("MILEXLEGAL  ·  AUXILIAR JURÍDICO CON IA", 14, DORADO, True, SANS)]])
texto(s, Inches(0.9), Inches(2.7), CW, Inches(1.6),
      [[R("Revisa tu Contrato", 54, CREMA, True, SERIF)]])
texto(s, Inches(0.9), Inches(4.3), Inches(9.6), Inches(1.6),
      [[R("Antes de firmar, entiende lo que firmas. La IA que hace el primer barrido de cualquier contrato en segundos — para que el abogado decida con todo sobre la mesa.", 17, CLARO, False, SANS)]])
texto(s, Inches(0.9), Inches(6.4), CW, Inches(0.5),
      [[R("hackathonlawgic.vercel.app", 14, DORADO, True, SANS),
        R("        Mini-Hackathon Lawgic 2026  ·  Brenda Barragán", 12, TENUE, False, SANS)]])

# ============ SLIDE 2: QUE ES ============
s = prs.slides.add_slide(BLANK)
encabezado(s, "¿QUÉ ES?", "Tu auxiliar jurídico con IA")
texto(s, Inches(1.15), Inches(2.3), CW, Inches(1.3),
      [[R("Una aplicación web que ", 16, TINTA),
        R("analiza y crea contratos legales", 16, NEGRO, True),
        R(" en español simple. En segundos detecta cláusulas riesgosas, las explica, sugiere qué preguntar, las relaciona con jurisprudencia aplicable y entrega una versión mejorada lista para Word.", 16, TINTA)]])
b = caja(s, Inches(1.15), Inches(4.0), CW - Inches(0.25), Inches(2.4), SOFT, LINEA, 1.0)
texto(s, Inches(1.5), Inches(4.25), CW - Inches(0.9), Inches(1.9),
      [[R("El problema que resuelve", 15, NEGRO, True, SANS)],
       [R("Revisar cada contrato cláusula por cláusula consume horas valiosas del despacho. MILEXLEGAL hace ese primer barrido al instante — para que el abogado dedique su tiempo al criterio legal, la estrategia y sus clientes.", 15, GRIS)]], sp_after=8)

# ============ SLIDE 3: PARA QUIEN ============
s = prs.slides.add_slide(BLANK)
encabezado(s, "¿PARA QUIÉN ES?", "Pensada para tres situaciones reales")
tarjetas = [
    ("⏱  Abogados con prisa",
     "Necesitan redactar algo por voz mientras hacen otra actividad. Dictan el contrato y la IA lo escribe completo, sin frenar su día."),
    ("⚖  Recién egresados (y no tanto)",
     "Aprenden a revisar contratos con seguridad: la IA señala riesgos y los explica en simple, como un mentor que adelanta el primer filtro."),
    ("$  Saber cuánto cobrar",
     "Una calculadora les dice cuánto cobrar por cada contrato según complejidad, servicios, urgencia y tipo de cliente."),
]
cw3 = (CW - Inches(0.6)) / 3
for i, (t, d) in enumerate(tarjetas):
    x = Inches(1.15) + i * (cw3 + Inches(0.3))
    caja(s, x, Inches(2.5), cw3, Inches(3.6), SOFT, LINEA, 1.0)
    texto(s, x + Inches(0.25), Inches(2.75), cw3 - Inches(0.5), Inches(0.9),
          [[R(t, 16, NEGRO, True, SANS)]])
    texto(s, x + Inches(0.25), Inches(3.7), cw3 - Inches(0.5), Inches(2.2),
          [[R(d, 13.5, GRIS)]])

# ============ SLIDE 4: SERVICIO 1 ============
s = prs.slides.add_slide(BLANK)
encabezado(s, "SERVICIO 01", "Analizamos tu contrato")
texto(s, Inches(1.15), Inches(2.15), CW, Inches(0.6),
      [[R("Subes un contrato en ", 14, TINTA), R("PDF, Word o texto", 14, NEGRO, True),
        R(" y la IA lo revisa cláusula por cláusula. Entrega un reporte completo:", 14, TINTA)]])
items = [
    ("\U0001F6A6  Nivel de riesgo", "Semáforo, puntaje de 0 a 100 y un veredicto que resume si es seguro firmar."),
    ("\U0001F534  Cláusulas riesgosas", "Cada cláusula problemática explicada en simple, con un “qué hacer”."),
    ("❓  Qué preguntar", "Las preguntas clave que debes hacer antes de firmar."),
    ("⚖  Jurisprudencia aplicable", "Tesis y referencias legales que aplican a ESE contrato (el plus)."),
    ("\U0001F6E1  Protecciones que faltan", "Cláusulas que el contrato debería tener y no tiene."),
    ("✅  Puntos a tu favor", "Lo que sí te protege en el contrato."),
    ("✍  Contrato mejorado + PDF", "Versión más justa lista para Word, y el análisis descargable en PDF."),
]
y = Inches(2.85)
for t, d in items:
    caja(s, Inches(1.15), y + Inches(0.06), Inches(0.15), Inches(0.15), DORADO)
    texto(s, Inches(1.5), y, CW - Inches(0.5), Inches(0.5),
          [[R(t + "  —  ", 13.5, NEGRO, True), R(d, 13.5, GRIS)]], sp_after=0)
    y = y + Inches(0.62)

# ============ SLIDE 5: SERVICIO 2 ============
s = prs.slides.add_slide(BLANK)
encabezado(s, "SERVICIO 02", "Creamos tu contrato — por voz")
texto(s, Inches(1.15), Inches(2.3), CW, Inches(1.2),
      [[R("Describes con tus palabras el contrato que necesitas — ", 16, TINTA),
        R("escribiendo o dictándolo por voz", 16, NEGRO, True),
        R(" — y la IA lo redacta completo, listo para descargar en Word y editar. Ideal para el abogado que va de prisa y dicta mientras hace otra cosa.", 16, TINTA)]])
b = caja(s, Inches(1.15), Inches(4.1), CW - Inches(0.25), Inches(2.3), SOFT, LINEA, 1.0)
texto(s, Inches(1.5), Inches(4.35), CW - Inches(0.9), Inches(1.9),
      [[R("\U0001F3A4  Conoce a “Lexi”, el asistente de voz", 15, NEGRO, True, SANS)],
       [R("Un muñequito animado que mueve la boca mientras hablas, muestra en vivo lo que vas dictando en una burbuja, y entra en modo “redactando” mientras la IA trabaja. Hace la experiencia más humana y fácil.", 14.5, GRIS)]], sp_after=8)

# ============ SLIDE 6: CALCULADORA ============
s = prs.slides.add_slide(BLANK)
encabezado(s, "HERRAMIENTA EXTRA", "Calculadora de precio")
texto(s, Inches(1.15), Inches(2.3), CW, Inches(1.0),
      [[R("Para que ningún abogado (sobre todo el recién egresado) dude cuánto cobrar: calcula al instante el precio de un contrato según varios factores.", 16, TINTA)]])
factores = [
    ("Complejidad", "Sencillo, estándar, complejo o corporativo"),
    ("Servicios", "Análisis, redacción, jurisprudencia, bilingue…"),
    ("Tamaño", "Páginas y número de revisiones incluidas"),
    ("Urgencia y cliente", "Normal/rápida/urgente · persona, PyME o corporativo"),
]
cw2 = (CW - Inches(0.3)) / 2
for i, (t, d) in enumerate(factores):
    col = i % 2; row = i // 2
    x = Inches(1.15) + col * (cw2 + Inches(0.3))
    yy = Inches(3.6) + row * Inches(1.5)
    caja(s, x, yy, cw2, Inches(1.25), SOFT, LINEA, 1.0)
    texto(s, x + Inches(0.25), yy + Inches(0.18), cw2 - Inches(0.5), Inches(1.0),
          [[R(t, 15, NEGRO, True, SANS)], [R(d, 13, GRIS)]], sp_after=3)

# ============ SLIDE 7: COMO FUNCIONA ============
s = prs.slides.add_slide(BLANK)
encabezado(s, "ASÍ DE FÁCIL", "Cómo funciona, en 3 pasos")
pasos = [
    ("1", "Sube o pega", "Tu contrato en PDF, Word o texto — o descríbelo por voz para crear uno nuevo."),
    ("2", "La IA trabaja", "Analiza cada cláusula y detecta riesgos, o redacta un contrato a tu medida."),
    ("3", "Descarga", "Reporte en PDF, y contratos mejorados o nuevos en Word, listos para usar."),
]
for i, (n, t, d) in enumerate(pasos):
    x = Inches(1.15) + i * (cw3 + Inches(0.3))
    caja(s, x, Inches(2.6), cw3, Inches(3.4), SOFT, LINEA, 1.0)
    texto(s, x + Inches(0.25), Inches(2.8), cw3 - Inches(0.5), Inches(0.9),
          [[R(n, 40, DORADO, True, SERIF)]])
    texto(s, x + Inches(0.25), Inches(3.8), cw3 - Inches(0.5), Inches(0.6),
          [[R(t, 17, NEGRO, True, SANS)]])
    texto(s, x + Inches(0.25), Inches(4.4), cw3 - Inches(0.5), Inches(1.5),
          [[R(d, 13.5, GRIS)]])

# ============ SLIDE 8: NUMEROS ============
s = prs.slides.add_slide(BLANK)
encabezado(s, "LO QUE OBTIENES", "Mucho valor en un solo lugar")
nums = [
    ("6", "revisiones en un solo análisis"),
    ("3", "formatos que lee: PDF, Word y texto"),
    ("2", "formas de exportar: PDF y Word"),
]
for i, (n, d) in enumerate(nums):
    x = Inches(1.15) + i * (cw3 + Inches(0.3))
    caja(s, x, Inches(2.8), cw3, Inches(2.8), NEGRO)
    texto(s, x, Inches(3.1), cw3, Inches(1.3),
          [[R(n, 66, DORADO2, True, SERIF)]], align=PP_ALIGN.CENTER)
    texto(s, x + Inches(0.3), Inches(4.5), cw3 - Inches(0.6), Inches(0.9),
          [[R(d, 14, CLARO)]], align=PP_ALIGN.CENTER)

# ============ SLIDE 9: FILOSOFIA ============
s = prs.slides.add_slide(BLANK)
fondo(s, NEGRO)
caja(s, Inches(0.9), Inches(2.4), Inches(1.1), Pt(4), DORADO)
texto(s, Inches(0.9), Inches(2.55), CW, Inches(0.5),
      [[R("\U0001F3AF  NUESTRA FILOSOFÍA", 14, DORADO, True, SANS)]])
texto(s, Inches(0.9), Inches(3.1), CW, Inches(1.4),
      [[R("Un auxiliar para el abogado, no un reemplazo.", 34, CREMA, True, SERIF)]])
texto(s, Inches(0.9), Inches(4.6), Inches(10.5), Inches(1.6),
      [[R("La herramienta apoya el trabajo y hace el primer barrido; el juicio profesional y la decisión final siempre son del abogado.", 18, CLARO)]])

# ============ SLIDE 10: TECNOLOGIA + PRUEBALO ============
s = prs.slides.add_slide(BLANK)
encabezado(s, "PRUÉBALO TÚ MISMO", "Está en vivo, en internet")
texto(s, Inches(1.15), Inches(2.4), CW, Inches(0.8),
      [[R("\U0001F517  ", 22, DORADO), R("hackathonlawgic.vercel.app", 24, NEGRO, True, SANS)]])
b = caja(s, Inches(1.15), Inches(3.4), CW - Inches(0.25), Inches(1.5), SOFT, LINEA, 1.0)
texto(s, Inches(1.5), Inches(3.62), CW - Inches(0.9), Inches(1.2),
      [[R("Sugerencia para la demo:", 14, NEGRO, True, SANS),
        R("  entra a Crear contrato → Dictar por voz (en Chrome o Edge) y habla — verás a Lexi mover la boca y escribir lo que dices. Luego ve a Analizar y prueba el ejemplo para ver el semáforo de riesgo.", 14, GRIS)]])
texto(s, Inches(1.15), Inches(5.3), CW, Inches(0.9),
      [[R("Hecha con  ", 13, GRIS), R("React + Vite", 13, TINTA, True),
        R("  ·  IA  ", 13, GRIS), R("Gemini (Google)", 13, TINTA, True),
        R("  ·  ", 13, GRIS), R("Supabase", 13, TINTA, True),
        R("  ·  ", 13, GRIS), R("Vercel", 13, TINTA, True),
        R("   — funciona en computadora y celular, modo claro y oscuro.", 13, GRIS)]])
texto(s, Inches(1.15), Inches(6.5), CW, Inches(0.7),
      [[R("MILEXLEGAL ofrece orientación informativa generada por IA y no sustituye la asesoría de un abogado.", 10.5, GRIS, False, SANS, True)]])

out = r"C:\Users\brenb\OneDrive\Desktop\Demo-MILEXLEGAL.pptx"
prs.save(out)
print("GUARDADO:", out, "-", len(prs.slides.__iter__.__self__._sldIdLst), "slides")

